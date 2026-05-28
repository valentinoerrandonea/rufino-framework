"""Format converters for process-batch staging."""
import logging
import shutil
import subprocess
from pathlib import Path

from rufino.engine.process.batch.errors import (
    ConversionError,
    UnsupportedFormatError,
)

_log = logging.getLogger(__name__)


class SofficeNotFoundError(RuntimeError):
    """soffice missing from PATH — system config problem, not a conversion failure."""


def _which_soffice() -> str | None:
    """Indirection so tests can monkeypatch the lookup."""
    return shutil.which("soffice")


def _cleanup_partial(out_pdf: Path) -> None:
    if out_pdf.exists():
        try:
            out_pdf.unlink()
        except OSError as e:
            _log.warning("could not remove partial PDF %s: %s", out_pdf, e)


def convert_to_pdf(
    path: Path, *, out_dir: Path, timeout: float = 120.0,
) -> Path:
    soffice = _which_soffice()
    if soffice is None:
        raise SofficeNotFoundError(
            "soffice not found in PATH — install LibreOffice "
            "(macOS: brew install --cask libreoffice) to use multimodal mode."
        )
    out_dir.mkdir(parents=True, exist_ok=True)
    out_pdf = out_dir / (path.stem + ".pdf")
    cmd = [
        soffice, "--headless", "--convert-to", "pdf",
        "--outdir", str(out_dir), str(path),
    ]
    try:
        result = subprocess.run(  # noqa: S603 (cmd built from trusted args)
            cmd, capture_output=True, timeout=timeout, check=False,
        )
    except subprocess.TimeoutExpired as e:
        _cleanup_partial(out_pdf)
        raise ConversionError(
            f"soffice conversion timed out after {timeout}s for {path}"
        ) from e
    if result.returncode != 0:
        _cleanup_partial(out_pdf)
        stderr = result.stderr.decode("utf-8", errors="replace")[:500]
        raise ConversionError(
            f"soffice failed for {path}: exit={result.returncode} "
            f"stderr={stderr}"
        )
    # soffice exit 0 is not a guarantee — historic versions emit empty or
    # truncated PDFs on certain DOCX inputs without flagging the failure.
    if not out_pdf.exists():
        raise ConversionError(
            f"soffice did not produce expected output {out_pdf}"
        )
    if out_pdf.stat().st_size == 0:
        _cleanup_partial(out_pdf)
        raise ConversionError(f"soffice produced empty PDF for {path}")
    if out_pdf.read_bytes()[:4] != b"%PDF":
        _cleanup_partial(out_pdf)
        raise ConversionError(f"soffice produced non-PDF output for {path}")
    if result.stderr:
        # User opted into multimodal specifically for fidelity; surface
        # soffice's non-fatal warnings (substituted font, missing image,
        # etc.) so they can spot lost diagrams instead of silently
        # accepting a degraded PDF.
        _log.warning(
            "soffice produced warnings (rc=0) for %s: %s",
            path, result.stderr.decode("utf-8", errors="replace")[:500],
        )
    return out_pdf


def docx_to_md(path: Path) -> str:
    """Convert a .docx file to markdown using mammoth.

    The ``mammoth`` dependency is imported lazily so the module itself remains
    importable in a stripped-down environment; only callers that actually need
    .docx conversion pay the cost (and see an error if the dep is missing).
    """
    try:
        import mammoth  # noqa: WPS433  (lazy import is intentional)
    except ImportError as e:
        raise RuntimeError(
            "mammoth is required for .docx conversion. "
            "Install it (it is declared in pyproject.toml) via "
            "`./install.sh` or `pipx install -e .` / `pip install -e .`."
        ) from e
    try:
        with open(path, "rb") as fh:
            result = mammoth.convert_to_markdown(fh)
        for msg in result.messages:
            text = getattr(msg, "message", None) or str(msg)
            _log.warning("docx conversion warning for %s: %s", path, text)
        return result.value
    except Exception as e:
        raise ConversionError(f"docx conversion failed for {path}: {e}") from e


def pptx_to_md(path: Path) -> str:
    """Convert a .pptx file to markdown, one '## Slide N: title' section per slide.

    The ``python-pptx`` dependency is imported lazily so the module itself
    remains importable when the dep is absent; only callers that actually need
    .pptx conversion pay the cost (and see an error if the dep is missing).
    """
    try:
        from pptx import Presentation  # noqa: WPS433  (lazy import is intentional)
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for .pptx conversion. "
            "Install it (it is declared in pyproject.toml) via "
            "`./install.sh` or `pipx install -e .` / `pip install -e .`."
        ) from e
    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise ConversionError(f"pptx open failed for {path}: {e}") from e

    sections: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        title = texts[0] if texts else f"(empty slide {idx})"
        body = "\n\n".join(texts[1:]) if len(texts) > 1 else ""
        section = f"## Slide {idx}: {title}\n"
        if body:
            section += f"\n{body}\n"
        sections.append(section)
    return "\n".join(sections)


def convert_to_markdown(path: Path) -> str:
    """Dispatch to the right converter (or passthrough) based on extension.

    Raises UnsupportedFormatError for unsupported extensions, legacy formats,
    and PDFs (which are passthrough at the stager level — callers should not
    invoke this for PDFs).
    """
    suffix = path.suffix.lower()
    if suffix in (".md", ".txt"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".docx":
        return docx_to_md(path)
    if suffix == ".pptx":
        return pptx_to_md(path)
    if suffix == ".pdf":
        raise UnsupportedFormatError(
            f"{suffix} is passthrough — copy the file verbatim; "
            f"do not call convert_to_markdown"
        )
    if suffix in (".doc", ".ppt"):
        raise UnsupportedFormatError(
            f"legacy {suffix} requires pandoc/LibreOffice — convert to "
            f"{suffix}x first (out of scope for v0.1.0)"
        )
    raise UnsupportedFormatError(f"{suffix} not supported")
