import logging
import shutil
from pathlib import Path
from types import SimpleNamespace

import pytest

pytest.importorskip("mammoth")
pytest.importorskip("pptx")

from rufino.engine.process.batch.converters import (
    SofficeNotFoundError,
    convert_to_markdown,
    convert_to_pdf,
    docx_to_md,
    pptx_to_md,
)
from rufino.engine.process.batch.errors import (
    ConversionError,
    UnsupportedFormatError,
)


FIXTURES = Path(__file__).parent / "fixtures" / "batch"


def test_converters_module_importable_without_mammoth(monkeypatch):
    """If mammoth is missing, importing converters must still succeed; only the
    docx call should raise a clear error."""
    import sys
    monkeypatch.setitem(sys.modules, "mammoth", None)
    monkeypatch.delitem(
        sys.modules, "rufino.engine.process.batch.converters", raising=False,
    )
    from rufino.engine.process.batch import converters  # must not raise
    with pytest.raises(RuntimeError, match="mammoth is required"):
        converters.convert_to_markdown(Path("x.docx"))


def test_converters_module_importable_without_pptx(monkeypatch):
    """If python-pptx is missing, importing converters must still succeed; only
    the pptx call should raise a clear error."""
    import sys
    monkeypatch.setitem(sys.modules, "pptx", None)
    monkeypatch.delitem(
        sys.modules, "rufino.engine.process.batch.converters", raising=False,
    )
    from rufino.engine.process.batch import converters  # must not raise
    with pytest.raises(RuntimeError, match="python-pptx is required"):
        converters.convert_to_markdown(Path("x.pptx"))


def test_docx_to_md_extracts_text():
    md = docx_to_md(FIXTURES / "hello.docx")
    assert "Hello from docx" in md


def test_pptx_to_md_extracts_per_slide():
    md = pptx_to_md(FIXTURES / "hello.pptx")
    assert "Slide one title" in md
    assert "Slide two title" in md
    assert md.count("## Slide") == 2


def test_convert_md_passthrough(tmp_path):
    src = tmp_path / "note.md"
    src.write_text("# already markdown\n")
    assert convert_to_markdown(src) == "# already markdown\n"


def test_convert_txt_passthrough(tmp_path):
    src = tmp_path / "note.txt"
    src.write_text("plain text\n")
    assert convert_to_markdown(src) == "plain text\n"


def test_convert_docx_dispatches():
    md = convert_to_markdown(FIXTURES / "hello.docx")
    assert "Hello from docx" in md


def test_convert_pptx_dispatches():
    md = convert_to_markdown(FIXTURES / "hello.pptx")
    assert "Slide one title" in md


def test_convert_legacy_doc_raises(tmp_path):
    legacy = tmp_path / "old.doc"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0")
    with pytest.raises(UnsupportedFormatError, match=r"\.doc"):
        convert_to_markdown(legacy)


def test_convert_legacy_ppt_raises(tmp_path):
    legacy = tmp_path / "old.ppt"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0")
    with pytest.raises(UnsupportedFormatError, match=r"\.ppt"):
        convert_to_markdown(legacy)


def test_convert_unknown_raises(tmp_path):
    unknown = tmp_path / "thing.xyz"
    unknown.write_text("?")
    with pytest.raises(UnsupportedFormatError, match=r"\.xyz"):
        convert_to_markdown(unknown)


def test_convert_pdf_does_not_call_converter(tmp_path):
    """PDFs are passthrough — converters should NOT be invoked. The caller
    leaves the PDF in place and lets the worker's Read tool handle it.
    """
    pdf = tmp_path / "doc.pdf"
    pdf.write_bytes(b"%PDF-1.4 ...")
    with pytest.raises(UnsupportedFormatError, match="passthrough"):
        convert_to_markdown(pdf)


def test_pptx_body_paragraphs_use_hard_break(tmp_path):
    """Body paragraphs in a pptx slide must be separated by a blank line so
    CommonMark renders them as distinct paragraphs (hard break), not a single
    paragraph with a soft line break.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout — full control over shapes
    slide = prs.slides.add_slide(blank)

    # Title textbox (first shape with text → becomes the slide title).
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text_frame.text = "The Title"

    # Body textbox with two paragraphs.
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = body_box.text_frame
    tf.text = "First paragraph"
    p2 = tf.add_paragraph()
    p2.text = "Second paragraph"

    out = tmp_path / "two_paragraphs.pptx"
    prs.save(str(out))

    md = pptx_to_md(out)

    assert "First paragraph" in md
    assert "Second paragraph" in md
    # The two body paragraphs must be separated by a blank line.
    assert "First paragraph\n\nSecond paragraph" in md
    # And the single-newline (soft break) form must NOT appear.
    assert "First paragraph\nSecond paragraph" not in md


def test_docx_warnings_are_logged(tmp_path, caplog, monkeypatch):
    """mammoth surfaces unconvertible elements via result.messages — those
    must be emitted as logger warnings, not silently dropped.
    """
    import mammoth

    fake_messages = [
        SimpleNamespace(type="warning", message="image dropped"),
        SimpleNamespace(type="warning", message="unrecognized style 'Foo'"),
    ]
    fake_result = SimpleNamespace(value="# converted body\n", messages=fake_messages)

    def fake_convert_to_markdown(fh):
        return fake_result

    monkeypatch.setattr(mammoth, "convert_to_markdown", fake_convert_to_markdown)

    src = tmp_path / "fake.docx"
    src.write_bytes(b"not a real docx but mammoth is stubbed")

    with caplog.at_level(logging.WARNING, logger="rufino.engine.process.batch.converters"):
        out = docx_to_md(src)

    assert out == "# converted body\n"
    messages = [rec.getMessage() for rec in caplog.records]
    assert any("image dropped" in m for m in messages)
    assert any("unrecognized style 'Foo'" in m for m in messages)


@pytest.mark.skipif(
    shutil.which("soffice") is None,
    reason="soffice not in PATH (LibreOffice not installed)",
)
def test_convert_to_pdf_produces_valid_pdf_for_docx(
    tmp_path: Path, batch_fixtures_dir: Path,
):
    src = batch_fixtures_dir / "hello.docx"
    out = convert_to_pdf(src, out_dir=tmp_path)
    assert out.exists()
    assert out.suffix == ".pdf"
    # PDFs start with `%PDF-`.
    assert out.read_bytes()[:5] == b"%PDF-"


@pytest.mark.skipif(
    shutil.which("soffice") is None,
    reason="soffice not in PATH (LibreOffice not installed)",
)
def test_convert_to_pdf_produces_valid_pdf_for_pptx(
    tmp_path: Path, batch_fixtures_dir: Path,
):
    src = batch_fixtures_dir / "hello.pptx"
    out = convert_to_pdf(src, out_dir=tmp_path)
    assert out.exists()
    assert out.read_bytes()[:5] == b"%PDF-"


def test_convert_to_pdf_raises_when_soffice_missing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
):
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: None,
    )
    src = tmp_path / "fake.docx"
    src.write_bytes(b"dummy")
    with pytest.raises(SofficeNotFoundError, match="soffice not found"):
        convert_to_pdf(src, out_dir=tmp_path)


def _fake_completed(returncode: int, stderr: bytes = b""):
    return SimpleNamespace(returncode=returncode, stderr=stderr, stdout=b"")


def test_convert_to_pdf_raises_on_timeout(tmp_path, monkeypatch):
    import subprocess

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _raise_timeout(*args, **kwargs):
        raise subprocess.TimeoutExpired(cmd=args[0], timeout=1.0)

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _raise_timeout,
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="timed out"):
        convert_to_pdf(src, out_dir=tmp_path, timeout=1.0)


def test_convert_to_pdf_raises_on_nonzero_exit(tmp_path, monkeypatch):
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run",
        lambda *a, **kw: _fake_completed(1, b"boom"),
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="exit=1"):
        convert_to_pdf(src, out_dir=tmp_path)


def test_convert_to_pdf_raises_when_output_missing(tmp_path, monkeypatch):
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run",
        lambda *a, **kw: _fake_completed(0),
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="did not produce"):
        convert_to_pdf(src, out_dir=tmp_path)


def test_convert_to_pdf_raises_on_empty_output(tmp_path, monkeypatch):
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _make_empty(*a, **kw):
        out = tmp_path / "doc.pdf"
        out.write_bytes(b"")
        return _fake_completed(0)

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _make_empty,
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="empty PDF"):
        convert_to_pdf(src, out_dir=tmp_path)
    assert not (tmp_path / "doc.pdf").exists()  # partial cleanup


def test_convert_to_pdf_raises_on_non_pdf_output(tmp_path, monkeypatch):
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _make_garbage(*a, **kw):
        out = tmp_path / "doc.pdf"
        out.write_bytes(b"not a pdf at all")
        return _fake_completed(0)

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _make_garbage,
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="non-PDF output"):
        convert_to_pdf(src, out_dir=tmp_path)
    assert not (tmp_path / "doc.pdf").exists()


def test_convert_to_pdf_warns_on_rc0_stderr(tmp_path, monkeypatch, caplog):
    """soffice rc=0 with non-empty stderr (font sub, dropped image) is
    surfaced as a warning, not silently swallowed — multimodal users
    opted in for fidelity, so they need the diagnostic."""
    import logging as _logging

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _make_real_pdf(*a, **kw):
        out = tmp_path / "doc.pdf"
        out.write_bytes(b"%PDF-1.4\nbody\n%%EOF")
        return _fake_completed(0, b"warning: substituted font Helvetica")

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _make_real_pdf,
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with caplog.at_level(
        _logging.WARNING, logger="rufino.engine.process.batch.converters",
    ):
        result = convert_to_pdf(src, out_dir=tmp_path)
    assert result.exists()
    assert any("substituted font" in r.getMessage() for r in caplog.records)


def test_cleanup_partial_swallows_oserror_with_warning(
    tmp_path, monkeypatch, caplog,
):
    """If cleanup itself can't unlink the partial PDF (read-only dir, etc.)
    the outer ConversionError still propagates and the cleanup failure is
    logged — not silently dropped."""
    import logging as _logging
    from pathlib import Path as _Path

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _make_empty(*a, **kw):
        out = tmp_path / "doc.pdf"
        out.write_bytes(b"")
        return _fake_completed(0)

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _make_empty,
    )

    real_unlink = _Path.unlink

    def _unlink_fails(self, *args, **kwargs):
        if self.name == "doc.pdf":
            raise OSError("read-only filesystem")
        return real_unlink(self, *args, **kwargs)

    monkeypatch.setattr(_Path, "unlink", _unlink_fails)

    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with caplog.at_level(
        _logging.WARNING, logger="rufino.engine.process.batch.converters",
    ):
        with pytest.raises(ConversionError, match="empty PDF"):
            convert_to_pdf(src, out_dir=tmp_path)
    assert any(
        "could not remove partial PDF" in r.getMessage() for r in caplog.records
    )


def test_convert_to_pdf_cleans_up_partial_on_nonzero_exit(tmp_path, monkeypatch):
    """soffice can write garbage then crash — leftover must not survive."""
    monkeypatch.setattr(
        "rufino.engine.process.batch.converters._which_soffice",
        lambda: "/fake/soffice",
    )

    def _write_then_fail(*a, **kw):
        out = tmp_path / "doc.pdf"
        out.write_bytes(b"partial...")
        return _fake_completed(1, b"crashed")

    monkeypatch.setattr(
        "rufino.engine.process.batch.converters.subprocess.run", _write_then_fail,
    )
    src = tmp_path / "doc.docx"
    src.write_bytes(b"x")
    with pytest.raises(ConversionError, match="exit=1"):
        convert_to_pdf(src, out_dir=tmp_path)
    assert not (tmp_path / "doc.pdf").exists()
